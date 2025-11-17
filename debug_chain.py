from server.services.registry import registry
import server.converters
from io import BytesIO
from pptx import Presentation
import json

print('Registry describe:')
print(json.dumps(registry.describe(), indent=2))

try:
    chain = registry.find_chain('pptx','pdf')
    print('Chain:', [(fn.__name__, ext) for fn, ext in chain])
except Exception as e:
    print('Chain error', e)

p = Presentation()
try:
    s = p.slides.add_slide(p.slide_layouts[5])
    tx = p.slides[0].shapes.title
    tx.text = 'Hello'
except Exception:
    pass
buf = BytesIO(); p.save(buf); buf.seek(0)
content = buf.read()

try:
    current = content
    mime = 'application/octet-stream'
    for func, next_ext in chain:
        print('Calling', func.__name__, '->', next_ext)
        current, mime = func(current, next_ext)
        print('Result mime:', mime, 'len', len(current))
    print('Final size', len(current), 'mime', mime)
except Exception as e:
    import traceback; traceback.print_exc()
