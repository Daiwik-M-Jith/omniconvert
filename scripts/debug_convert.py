from io import BytesIO
from pptx import Presentation
from fastapi.testclient import TestClient
from server.app import app


def create_sample_pptx():
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[5])
    try:
        txBox = pres.slides[0].shapes.title
        txBox.text = "Hello chain"
    except Exception:
        pass
    buf = BytesIO()
    pres.save(buf)
    buf.seek(0)
    return buf


def main():
    client = TestClient(app)
    buf = create_sample_pptx()
    files = {"file": ("sample.pptx", buf.read(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
    data = {"target_format": "pdf"}
    resp = client.post('/api/convert', data=data, files=files)
    print('status:', resp.status_code)
    print('headers:', resp.headers)
    try:
        print('body text:', resp.text[:200])
    except Exception:
        print('<no printable body>')


if __name__ == '__main__':
    main()
