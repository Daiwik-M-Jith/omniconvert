from io import BytesIO
from typing import Tuple

from openpyxl import load_workbook
from pptx import Presentation

from .base import register_converter


@register_converter("xlsx", "csv", note="Extracts first sheet to CSV via openpyxl")
def xlsx_to_csv(content: bytes, target: str = "csv") -> Tuple[bytes, str]:
    wb = load_workbook(filename=BytesIO(content), data_only=True)
    sheet = wb.active
    output = BytesIO()
    # simple CSV handling
    for row in sheet.iter_rows(values_only=True):
        row_text = ",".join(["" if v is None else str(v) for v in row])
        output.write((row_text + "\n").encode("utf-8"))
    return output.getvalue(), "text/csv"


@register_converter("pptx", "txt", note="Extract slides text from PPTX into plain text")
def pptx_to_txt(content: bytes, target: str = "txt") -> Tuple[bytes, str]:
    pres = Presentation(BytesIO(content))
    lines = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                lines.append(shape.text)
    return ("\n\n".join(lines)).encode("utf-8"), "text/plain"
