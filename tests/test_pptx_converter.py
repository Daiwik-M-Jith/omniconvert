from io import BytesIO

from pptx import Presentation

from server.services.registry import registry
import server.converters  # noqa: F401


def test_pptx_to_txt_converter_direct_call():
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[5])
    try:
        pres.slides[0].shapes.title.text = "Hello world"
    except Exception:
        pass
    buf = BytesIO(); pres.save(buf); buf.seek(0)
    converter = registry.resolve('pptx','txt')
    output, mime = converter(buf.read(), 'txt')
    assert isinstance(output, (bytes, bytearray))
    assert 'text' in mime
