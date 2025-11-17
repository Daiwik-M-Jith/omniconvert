from io import BytesIO

from openpyxl import Workbook
from pptx import Presentation
from fastapi.testclient import TestClient

from server.app import app
import server.converters  # noqa: F401

client = TestClient(app)


def test_xlsx_to_csv_conversion():
    wb = Workbook()
    ws = wb.active
    ws.append(["Name", "Age"])
    ws.append(["Alice", 30])
    buffer = BytesIO()
    wb.save(buffer)
    buffer.seek(0)

    files = {"file": ("sample.xlsx", buffer.read(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")}
    data = {"target_format": "csv"}
    response = client.post("/api/convert", data=data, files=files)
    assert response.status_code == 200
    assert response.headers["content-type"].startswith("text/csv")
    assert b"Alice" in response.content


def test_pptx_to_txt_conversion():
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[5])
    txBox = slide.shapes.title
    txBox.text = "Hello from Slide"
    buffer = BytesIO()
    pres.save(buffer)
    buffer.seek(0)

    files = {"file": ("sample.pptx", buffer.read(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
    data = {"target_format": "txt"}
    response = client.post("/api/convert", data=data, files=files)
    assert response.status_code == 200
    assert response.headers["content-type"].startswith("text/")
    assert b"Hello from Slide" in response.content
