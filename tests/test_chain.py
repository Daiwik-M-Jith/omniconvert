from io import BytesIO
from pptx import Presentation
from fastapi.testclient import TestClient

from server.app import app
import server.converters  # noqa: F401

client = TestClient(app)


def test_pptx_to_pdf_chain():
    # PPTX -> TXT -> PDF chain (pptx->txt available and txt->pdf available)
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[5])
    if pres.slides and hasattr(pres.slides[0], 'shapes') and pres.slides[0].shapes:
        pass
    # Add title per previous test approach
    try:
        txBox = pres.slides[0].shapes.title
        txBox.text = "Hello chain"
    except Exception:
        # ensure something on slide
        pass
    buffer = BytesIO()
    pres.save(buffer)
    buffer.seek(0)

    files = {"file": ("sample.pptx", buffer.read(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
    data = {"target_format": "pdf"}
    response = client.post("/api/convert", data=data, files=files)
    assert response.status_code == 200
    # PDF should begin with %PDF
    assert response.content.startswith(b"%PDF")
